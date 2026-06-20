"""Add study notes to PowerPoint speaker notes with Presenter View–friendly formatting."""
from __future__ import annotations

import re
import shutil
import sys
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from abbreviations import expand_abbreviations

ROOT = Path(__file__).resolve().parent
SLIDE_RE = re.compile(r"^## Slide (\d+) — (.+)$", re.M)
IMG_RE = re.compile(r'<img\b[^>]*\balt="([^"]*)"[^>]*>', re.I)

FONT_TITLE = Pt(14)
FONT_HEADING = Pt(12)
FONT_BODY = Pt(11)


@dataclass
class NoteBlock:
    text: str
    kind: str  # title, heading, subheading, bullet, text, quote, diagram, spacer
    level: int = 0


def parse_study_notes(md: str) -> dict[int, tuple[str, str]]:
    md = md.split("## Reference — Additional Topics")[0]
    notes: dict[int, tuple[str, str]] = {}
    for chunk in re.split(r"^---\s*$", md, flags=re.M):
        m = SLIDE_RE.search(chunk)
        if not m:
            continue
        num = int(m.group(1))
        title = m.group(2).strip()
        body = chunk[m.end() :].strip()
        notes[num] = (title, body)
    return notes


def clean_inline(text: str) -> str:
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return expand_abbreviations(text.strip())


def preprocess_body(body: str) -> str:
    body = IMG_RE.sub(r"\n[Diagram: \1]\n", body)
    body = re.sub(r"<img\b[^>]*>", "\n[Diagram]\n", body, flags=re.I)
    return body


def body_to_blocks(body: str) -> list[NoteBlock]:
    body = preprocess_body(body)
    blocks: list[NoteBlock] = []
    in_code = False

    for raw in body.splitlines():
        line = raw.rstrip()
        stripped = line.strip()

        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            blocks.append(NoteBlock(clean_inline(stripped), "text", level=1))
            continue

        if not stripped:
            if blocks and blocks[-1].kind != "spacer":
                blocks.append(NoteBlock("", "spacer"))
            continue

        if stripped.startswith("### "):
            blocks.append(NoteBlock(clean_inline(stripped[4:]), "heading"))
            continue
        if stripped.startswith("#### "):
            blocks.append(NoteBlock(clean_inline(stripped[5:]), "subheading"))
            continue
        if stripped.startswith("## "):
            blocks.append(NoteBlock(clean_inline(stripped[3:]), "heading"))
            continue

        if stripped.startswith("[Diagram"):
            blocks.append(NoteBlock(clean_inline(stripped), "diagram"))
            continue

        if re.match(r"^[-*]\s+", stripped):
            blocks.append(NoteBlock(clean_inline(stripped[2:]), "bullet", level=1))
            continue

        if stripped.startswith(">"):
            blocks.append(NoteBlock(clean_inline(stripped.lstrip("> ")), "quote", level=1))
            continue

        if stripped.startswith("|"):
            if re.match(r"^\|?\s*:?-+", stripped):
                continue
            cells = [clean_inline(c.strip()) for c in stripped.strip("|").split("|") if c.strip()]
            if cells:
                blocks.append(NoteBlock("  • " + " | ".join(cells), "text", level=1))
            continue

        blocks.append(NoteBlock(clean_inline(stripped), "text"))

    return blocks


def style_paragraph(paragraph, block: NoteBlock) -> None:
    if block.kind == "spacer":
        paragraph.text = ""
        return

    paragraph.text = block.text
    paragraph.level = block.level

    font = paragraph.font
    font.name = "Calibri"
    font.size = FONT_BODY
    font.bold = False
    font.italic = False

    if block.kind == "title":
        font.bold = True
        font.size = FONT_TITLE
    elif block.kind == "heading":
        font.bold = True
        font.size = FONT_HEADING
    elif block.kind == "subheading":
        font.bold = True
        font.size = FONT_BODY
    elif block.kind == "quote":
        font.italic = True
        paragraph.level = max(paragraph.level, 1)
    elif block.kind == "diagram":
        font.italic = True
    elif block.kind == "bullet":
        paragraph.level = max(paragraph.level, 1)


def set_slide_notes(slide, slide_num: int, title: str, body: str) -> None:
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()

    blocks: list[NoteBlock] = [
        NoteBlock(expand_abbreviations(f"Slide {slide_num} — {title}"), "title"),
        NoteBlock("", "spacer"),
        *body_to_blocks(body),
    ]

    first = True
    for block in blocks:
        if first:
            paragraph = notes_frame.paragraphs[0]
            first = False
        else:
            paragraph = notes_frame.add_paragraph()
        style_paragraph(paragraph, block)


def apply_notes(
    pptx_path: Path,
    study_path: Path,
    output_path: Path | None = None,
    backup: bool = True,
) -> tuple[int, int]:
    if not pptx_path.exists():
        raise FileNotFoundError(pptx_path)
    if not study_path.exists():
        raise FileNotFoundError(study_path)

    out = output_path or pptx_path
    if backup and out == pptx_path:
        bak = pptx_path.with_suffix(pptx_path.suffix + ".bak")
        if not bak.exists():
            shutil.copy2(pptx_path, bak)

    slide_notes = parse_study_notes(study_path.read_text(encoding="utf-8"))
    prs = Presentation(str(pptx_path))
    applied = 0
    for num, (title, body) in sorted(slide_notes.items()):
        idx = num - 1
        if idx < 0 or idx >= len(prs.slides):
            continue
        set_slide_notes(prs.slides[idx], num, title, body)
        applied += 1

    prs.save(str(out))
    return applied, len(prs.slides)


def main() -> None:
    repo_root = ROOT.parents[1]
    default_study = ROOT / "day1_study_notes.md"

    if len(sys.argv) > 1:
        pptx = Path(sys.argv[1])
    else:
        pptx = repo_root.parent / "Cybersecurity-AI-day1.pptx"

    if len(sys.argv) > 2:
        out_pptx = Path(sys.argv[2])
    else:
        out_pptx = pptx.with_name(pptx.stem + "-notes.pptx")

    if len(sys.argv) > 3:
        study_path = Path(sys.argv[3])
    else:
        # Infer study notes from deck name (day1 / day2)
        stem = pptx.stem.lower()
        if "day2" in stem:
            study_path = ROOT / "day2_study_notes.md"
        elif "day1" in stem:
            study_path = ROOT / "day1_study_notes.md"
        else:
            study_path = default_study

    applied, total = apply_notes(pptx, study_path, output_path=out_pptx, backup=(out_pptx == pptx))
    print(f"Saved {out_pptx}")
    print(f"  Study notes: {study_path.name}")
    print(f"  Applied formatted notes to {applied} slides (deck has {total} slides)")


if __name__ == "__main__":
    main()
